"""
Slide Generator — Create PowerPoint presentations from structured data.
Supports business presentations, pitches, reports, and custom decks.
"""

import os
import logging
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("kaihara.slide_generator")

# Output directory
OUTPUT_DIR = Path(os.getenv("KAIHARA_OUTPUT_DIR", "outputs/slides"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Color palette
COLORS = {
    'primary': RGBColor(0x8b, 0x5c, 0xf6),  # Purple
    'accent': RGBColor(0xa7, 0x8b, 0xfa),   # Light purple
    'dark': RGBColor(0x1a, 0x1a, 0x1a),
    'light': RGBColor(0xf5, 0xf5, 0xf5),
    'white': RGBColor(0xff, 0xff, 0xff),
    'gray': RGBColor(0x73, 0x73, 0x73),
}


def create_presentation(
    title: str,
    slides: list[dict],
    output_filename: Optional[str] = None,
    template: str = "default",
) -> str:
    """Create a PowerPoint presentation.
    
    Args:
        title: Presentation title
        slides: List of slide data, each with 'type' and content
                Types: 'title', 'content', 'two_column', 'image', 'divider'
        output_filename: Custom filename (without .pptx)
        template: Template style (default, dark, minimal)
    
    Returns:
        Path to generated .pptx file
    """
    if not output_filename:
        output_filename = f"presentation_{title.lower().replace(' ', '_')}"
    
    filepath = OUTPUT_DIR / f"{output_filename}.pptx"
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            _add_title_slide(prs, slide_data)
        elif slide_type == 'content':
            _add_content_slide(prs, slide_data)
        elif slide_type == 'two_column':
            _add_two_column_slide(prs, slide_data)
        elif slide_type == 'image':
            _add_image_slide(prs, slide_data)
        elif slide_type == 'divider':
            _add_divider_slide(prs, slide_data)
        elif slide_type == 'bullet':
            _add_bullet_slide(prs, slide_data)
    
    prs.save(str(filepath))
    
    logger.info(f"Presentation generated: {filepath}")
    return str(filepath)


def _add_title_slide(prs: Presentation, data: dict):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = data.get('title', 'Presentation')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if data.get('subtitle'):
        p2 = tf.add_paragraph()
        p2.text = data['subtitle']
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['light']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    
    # Author/Date
    if data.get('author'):
        author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
        tf2 = author_box.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = data['author']
        p3.font.size = Pt(16)
        p3.font.color.rgb = COLORS['light']
        p3.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, data: dict):
    """Add a content slide with title and body."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    content = data.get('content', '')
    if isinstance(content, list):
        for i, item in enumerate(content):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(10)
    else:
        p = tf2.paragraphs[0]
        p.text = content
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark']


def _add_two_column_slide(prs: Presentation, data: dict):
    """Add a two-column slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    left_content = data.get('left', '')
    if isinstance(left_content, list):
        for i, item in enumerate(left_content):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(8)
    else:
        p = tf_left.paragraphs[0]
        p.text = left_content
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    right_content = data.get('right', '')
    if isinstance(right_content, list):
        for i, item in enumerate(right_content):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(8)
    else:
        p = tf_right.paragraphs[0]
        p.text = right_content
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']


def _add_image_slide(prs: Presentation, data: dict):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    if data.get('title'):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['title']
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
    
    # Image
    image_path = data.get('image_path', '')
    if image_path and os.path.exists(image_path):
        left = (13.333 - 10) / 2
        slide.shapes.add_picture(image_path, Inches(left), Inches(1.5), Inches(10), Inches(5))


def _add_divider_slide(prs: Presentation, data: dict):
    """Add a divider/section slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']
    
    # Section number
    if data.get('number'):
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['number']
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(2))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = data.get('title', '')
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER


def _add_bullet_slide(prs: Presentation, data: dict):
    """Add a bullet point slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Bullets
    bullets_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5.5))
    tf2 = bullets_box.text_frame
    tf2.word_wrap = True
    
    bullets = data.get('bullets', [])
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(15)


def generate_business_pitch(
    company_name: str,
    problem: str,
    solution: str,
    market_size: str,
    business_model: str,
    team: list[str],
    ask: str,
    output_filename: Optional[str] = None,
) -> str:
    """Generate a business pitch deck.
    
    Args:
        company_name: Company name
        problem: Problem statement
        solution: Solution description
        market_size: Market size
        business_model: Business model
        team: Team members
        ask: What you're asking for
    
    Returns:
        Path to generated .pptx file
    """
    slides = [
        {'type': 'title', 'title': company_name, 'subtitle': 'Business Pitch'},
        {'type': 'bullet', 'title': 'The Problem', 'bullets': [problem]},
        {'type': 'bullet', 'title': 'Our Solution', 'bullets': [solution]},
        {'type': 'bullet', 'title': 'Market Opportunity', 'bullets': [market_size]},
        {'type': 'bullet', 'title': 'Business Model', 'bullets': [business_model]},
        {'type': 'bullet', 'title': 'Our Team', 'bullets': team},
        {'type': 'bullet', 'title': 'The Ask', 'bullets': [ask]},
    ]
    
    return create_presentation(company_name, slides, output_filename)


def generate_report_slides(
    title: str,
    sections: list[dict],
    output_filename: Optional[str] = None,
) -> str:
    """Generate a report presentation.
    
    Args:
        title: Report title
        sections: List of sections with 'title' and 'content'
    
    Returns:
        Path to generated .pptx file
    """
    slides = [
        {'type': 'title', 'title': title},
    ]
    
    for i, section in enumerate(sections):
        # Divider slide
        slides.append({
            'type': 'divider',
            'number': str(i + 1),
            'title': section.get('title', ''),
        })
        
        # Content slide
        content = section.get('content', '')
        if isinstance(content, list):
            slides.append({
                'type': 'bullet',
                'title': section.get('title', ''),
                'bullets': content,
            })
        else:
            slides.append({
                'type': 'content',
                'title': section.get('title', ''),
                'content': content,
            })
    
    return create_presentation(title, slides, output_filename)


# Tool registration for Kaihara
TOOLS = {
    "create_presentation": create_presentation,
    "generate_business_pitch": generate_business_pitch,
    "generate_report_slides": generate_report_slides,
}
